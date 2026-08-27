"""
Builds the TalentBridge AI project-submission deck as a real .pptx.

    python scripts/build-presentation.py [output_path]

Default output: %USERPROFILE%/Desktop/TalentBridge-AI-Overview.pptx

Everything asserted in these slides is taken from the codebase as it stands —
the scoring bands from src/lib/scoring.ts, the extraction contract from
src/lib/openai.ts, the pipeline stages from src/lib/pipeline.ts, the auth
methods from src/lib/auth.ts, and the counts on the Results slide from the
repository itself. Re-run this script after changing the app so the deck and
the implementation don't drift apart.

Fonts are deliberately Segoe UI / Consolas rather than the app's Space
Grotesk / Inter / IBM Plex Mono: PowerPoint silently substitutes fonts that
aren't installed on the opening machine, and the substitution is uglier than
picking a face that is always present on Windows.
"""

import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# --- Design tokens, lifted from tailwind.config.ts -------------------------
INK = RGBColor.from_string("16213D")
CANVAS = RGBColor.from_string("F6F5F2")
SURFACE = RGBColor.from_string("FFFFFF")
BORDER = RGBColor.from_string("DEDCD3")
ACCENT = RGBColor.from_string("1F6F6F")
ACCENT_SOFT = RGBColor.from_string("E4EFEF")
SUCCESS = RGBColor.from_string("2F7A4F")
SUCCESS_SOFT = RGBColor.from_string("E6F1E9")
WARNING = RGBColor.from_string("B07D22")
WARNING_SOFT = RGBColor.from_string("F6EDDC")
DANGER = RGBColor.from_string("B0433F")
DANGER_SOFT = RGBColor.from_string("F6E4E3")
MUTED = RGBColor.from_string("6B7280")
WHITE = RGBColor.from_string("FFFFFF")

FONT = "Segoe UI"
MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.85)
CONTENT_W = SLIDE_W - 2 * MARGIN

FOOTER_LEFT = "TalentBridge AI — Project Report"


# --- Primitives ------------------------------------------------------------
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def paint(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, left, top, width, height, fill=None, line=None, line_pt=1.0,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(shape, left, top, width, height)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_pt)
    sh.text_frame.word_wrap = True
    return sh


def text(slide, left, top, width, height, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.paragraphs[0].alignment = align
    return tf


def para(tf, first=False):
    """First paragraph of a fresh text frame is reused; later ones appended."""
    return tf.paragraphs[0] if first else tf.add_paragraph()


def run(p, s, size=14, color=INK, bold=False, italic=False, font=FONT):
    r = p.add_run()
    r.text = s
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return r


def line_of(tf, s, size=14, color=INK, bold=False, italic=False, font=FONT,
            first=False, space_after=8, space_before=0, align=None,
            indent_level=0):
    p = para(tf, first)
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.level = indent_level
    if align is not None:
        p.alignment = align
    run(p, s, size=size, color=color, bold=bold, italic=italic, font=font)
    return p


# --- Slide chrome ----------------------------------------------------------
def content_slide(prs, kicker, title, subtitle=None):
    slide = blank(prs)
    paint(slide, CANVAS)

    tf = text(slide, MARGIN, Inches(0.42), CONTENT_W, Inches(0.3))
    line_of(tf, kicker.upper(), size=11, color=ACCENT, bold=True, first=True,
            space_after=0)

    tf = text(slide, MARGIN, Inches(0.72), CONTENT_W, Inches(0.55))
    line_of(tf, title, size=30, color=INK, bold=True, first=True, space_after=0)

    box(slide, MARGIN, Inches(1.36), Inches(1.5), Inches(0.055), fill=ACCENT,
        shape=MSO_SHAPE.RECTANGLE)

    top = Inches(1.72)
    if subtitle:
        tf = text(slide, MARGIN, Inches(1.58), CONTENT_W, Inches(0.4))
        line_of(tf, subtitle, size=14, color=MUTED, first=True, space_after=0)
        top = Inches(2.12)
    return slide, top


def footer(slide, number):
    tf = text(slide, MARGIN, Inches(6.94), Inches(6.0), Inches(0.3))
    line_of(tf, FOOTER_LEFT, size=9, color=MUTED, first=True, space_after=0)

    tf = text(slide, SLIDE_W - MARGIN - Inches(1.0), Inches(6.94), Inches(1.0),
              Inches(0.3), align=PP_ALIGN.RIGHT)
    line_of(tf, str(number), size=9, color=MUTED, first=True, space_after=0,
            align=PP_ALIGN.RIGHT, font=MONO)


def notes(slide, s):
    slide.notes_slide.notes_text_frame.text = s


# --- Composite blocks ------------------------------------------------------
def bullets(slide, top, items, size=15, gap=13, width=None, left=None,
            bullet_color=ACCENT):
    """items: str | (str, str) for lead-in + rest | ('--', str) for a sub-item."""
    left = MARGIN if left is None else left
    width = (CONTENT_W if width is None else width)
    tf = text(slide, left, top, width, SLIDE_H - top - Inches(0.7))
    first = True
    for item in items:
        sub = isinstance(item, tuple) and item[0] == "--"
        p = para(tf, first)
        first = False
        p.space_after = Pt(gap if not sub else max(4, gap - 5))
        if sub:
            run(p, "      –   ", size=size - 1, color=MUTED)
            run(p, item[1], size=size - 1, color=MUTED)
            continue
        run(p, "■   ", size=size - 3, color=bullet_color, bold=True)
        if isinstance(item, tuple):
            run(p, item[0], size=size, color=INK, bold=True)
            run(p, item[1], size=size, color=INK)
        else:
            run(p, item, size=size, color=INK)
    return tf


def cards(slide, top, items, cols=2, height=Inches(1.55), gap=Inches(0.28),
          heading_size=15, body_size=12, accent=ACCENT, fill=SURFACE):
    """items: list of (heading, body)."""
    w = int((CONTENT_W - gap * (cols - 1)) / cols)
    for i, (heading, body) in enumerate(items):
        r, c = divmod(i, cols)
        left = MARGIN + c * (w + gap)
        t = top + r * (height + gap)
        sh = box(slide, left, t, w, height, fill=fill, line=BORDER)
        box(slide, left, t, Inches(0.05), height, fill=accent,
            shape=MSO_SHAPE.RECTANGLE)
        tf = sh.text_frame
        tf.margin_left = Inches(0.28)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.17)
        tf.margin_bottom = Inches(0.14)
        line_of(tf, heading, size=heading_size, color=INK, bold=True,
                first=True, space_after=5)
        if body:
            line_of(tf, body, size=body_size, color=MUTED, space_after=0)


def two_col(slide, top, left_title, left_items, right_title, right_items,
            left_accent=SUCCESS, right_accent=DANGER, left_fill=SUCCESS_SOFT,
            right_fill=DANGER_SOFT, height=Inches(4.15), item_size=12,
            item_space=7):
    w = int((CONTENT_W - Inches(0.35)) / 2)
    for i, (title, items, accent, fill) in enumerate([
        (left_title, left_items, left_accent, left_fill),
        (right_title, right_items, right_accent, right_fill),
    ]):
        left = MARGIN + i * (w + Inches(0.35))
        sh = box(slide, left, top, w, height, fill=fill, line=BORDER)
        tf = sh.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.22)
        line_of(tf, title, size=15, color=accent, bold=True, first=True,
                space_after=10)
        for it in items:
            p = para(tf)
            p.space_after = Pt(item_space)
            run(p, "•   ", size=item_size, color=accent, bold=True)
            run(p, it, size=item_size, color=INK)


def flow(slide, top, steps, height=Inches(1.05), arrow=Inches(0.3)):
    """Horizontal chevron flow. steps: list of (label, sublabel).

    A sublabel may contain "\\n" — it becomes a second paragraph rather than a
    literal newline inside a run, which python-pptx would escape as a control
    character.
    """
    n = len(steps)
    w = int((CONTENT_W - arrow * (n - 1)) / n)
    for i, (label, sub) in enumerate(steps):
        left = MARGIN + i * (w + arrow)
        sh = box(slide, left, top, w, height, fill=SURFACE, line=ACCENT,
                 line_pt=1.25)
        tf = sh.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.1)
        line_of(tf, label, size=12.5, color=INK, bold=True, first=True,
                space_after=3, align=PP_ALIGN.CENTER)
        for j, part in enumerate((sub or "").split("\n")):
            if not part:
                continue
            line_of(tf, part, size=10, color=MUTED, space_after=0,
                    align=PP_ALIGN.CENTER)
        if i < n - 1:
            a = box(slide, left + w, top + int(height / 2) - Inches(0.13),
                    arrow, Inches(0.26), fill=ACCENT,
                    shape=MSO_SHAPE.RIGHT_ARROW)
            a.line.fill.background()


def tiles(slide, top, stats, cols=4, height=Inches(1.3), gap=Inches(0.25)):
    """stats: list of (value, label)."""
    w = int((CONTENT_W - gap * (cols - 1)) / cols)
    for i, (value, label) in enumerate(stats):
        r, c = divmod(i, cols)
        left = MARGIN + c * (w + gap)
        t = top + r * (height + gap)
        sh = box(slide, left, t, w, height, fill=SURFACE, line=BORDER)
        tf = sh.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.22)
        line_of(tf, value, size=27, color=ACCENT, bold=True, first=True,
                space_after=2, font=MONO)
        line_of(tf, label, size=11, color=MUTED, space_after=0)


def table(slide, top, headers, rows, widths, row_h=Inches(0.46),
          head_fill=INK, size=12):
    xs, x = [], MARGIN
    for fr in widths:
        xs.append(x)
        x += int(CONTENT_W * fr)
    ws = [int(CONTENT_W * fr) for fr in widths]

    for i, h in enumerate(headers):
        sh = box(slide, xs[i], top, ws[i], row_h, fill=head_fill,
                 shape=MSO_SHAPE.RECTANGLE)
        tf = sh.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.15)
        line_of(tf, h, size=size, color=WHITE, bold=True, first=True,
                space_after=0)

    for r, row in enumerate(rows):
        t = top + row_h * (r + 1)
        fill = SURFACE if r % 2 == 0 else CANVAS
        for i, cell in enumerate(row):
            color, bold, font = INK, False, FONT
            if isinstance(cell, tuple):
                cell, color = cell[0], cell[1]
                bold = True
            if i == 0:
                bold = True
            sh = box(slide, xs[i], t, ws[i], row_h, fill=fill, line=BORDER,
                     line_pt=0.5, shape=MSO_SHAPE.RECTANGLE)
            tf = sh.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.15)
            tf.margin_right = Inches(0.1)
            line_of(tf, cell, size=size, color=color, bold=bold, first=True,
                    space_after=0, font=font)


# =========================================================================
#  Slides
# =========================================================================
def slide_title(prs):
    slide = blank(prs)
    paint(slide, INK)

    box(slide, Inches(0), Inches(0), Inches(0.28), SLIDE_H, fill=ACCENT,
        shape=MSO_SHAPE.RECTANGLE)

    tf = text(slide, Inches(1.15), Inches(1.55), Inches(10.6), Inches(0.4))
    line_of(tf, "PROJECT REPORT & SYSTEM DEMONSTRATION", size=12,
            color=ACCENT_SOFT, bold=True, first=True, space_after=0)

    tf = text(slide, Inches(1.15), Inches(2.12), Inches(10.6), Inches(1.1))
    line_of(tf, "TalentBridge AI", size=52, color=WHITE, bold=True,
            first=True, space_after=6)

    tf = text(slide, Inches(1.15), Inches(3.18), Inches(9.6), Inches(1.0))
    line_of(tf, "An evidence-based AI recruitment assistant that qualifies "
                "candidates against explicit job requirements — and leaves "
                "the hiring decision to a human.",
            size=17, color=RGBColor.from_string("C8D2DC"), first=True,
            space_after=0)

    box(slide, Inches(1.15), Inches(4.42), Inches(1.5), Inches(0.055),
        fill=ACCENT, shape=MSO_SHAPE.RECTANGLE)

    tf = text(slide, Inches(1.15), Inches(4.78), Inches(10.6), Inches(1.2))
    line_of(tf, "Gold Dick", size=15, color=WHITE, bold=True, first=True,
            space_after=5)
    line_of(tf, "[ Course · Module code · Supervisor · Institution ]",
            size=13, color=RGBColor.from_string("8A97A6"), space_after=5)
    line_of(tf, "August 2026", size=13,
            color=RGBColor.from_string("8A97A6"), space_after=0, font=MONO)

    tf = text(slide, Inches(1.15), Inches(6.62), Inches(10.6), Inches(0.4))
    line_of(tf, "Qualify first. Hire smarter.", size=13, color=ACCENT_SOFT,
            italic=True, first=True, space_after=0)

    notes(slide, "Opening. Replace the bracketed line with your course, "
                 "module code, supervisor and institution before submitting.\n\n"
                 "One-sentence framing: this project builds a recruitment "
                 "screening system where the AI extracts and explains "
                 "evidence, but a deterministic rule engine produces the "
                 "score and a human makes every decision.")
    return slide


def slide_contents(prs, n):
    slide, top = content_slide(prs, "Contents", "What this presentation covers")
    items = [
        ("1  Background & problem", "  — why CV screening needs more than a ranked list"),
        ("2  Aim & objectives", "  — what the system had to achieve (O1–O6)"),
        ("3  Scope & delimitations", "  — what was built, and what was deliberately excluded"),
        ("4  Methodology & design principles", "  — the four rules that shaped the architecture"),
        ("5  System architecture & data model", "  — layers, stack, and the 15-model schema"),
        ("6  Implementation", "  — pipeline, scoring engine, extraction contract, auth & access control"),
        ("7  Results", "  — what was delivered, mapped back to the objectives"),
        ("8  Evaluation & verification", "  — how correctness was established, and its limits"),
        ("9  Limitations & future work", "  — known gaps and the roadmap that follows from them"),
        ("10  Conclusion", ""),
    ]
    bullets(slide, top, items, size=14.5, gap=11)
    footer(slide, n)
    notes(slide, "Keep this brief — 20 seconds. Signpost that Results and "
                 "Limitations are where the honest assessment lives.")
    return slide


def slide_background(prs, n):
    slide, top = content_slide(
        prs, "1 · Background",
        "Screening is the bottleneck in hiring",
        "A single vacancy can attract hundreds of applications. The tools that "
        "exist to triage them trade away either accuracy or accountability.")
    cards(slide, top, [
        ("Manual screening does not scale",
         "Reading every CV against every requirement is linear work. Under time "
         "pressure it degrades into skimming the first page — so qualified "
         "candidates are missed for reasons nobody can reconstruct afterwards."),
        ("Keyword filters reject on absence, not evidence",
         "Conventional applicant-tracking filters match strings. A CV that says "
         "\"container orchestration in production\" fails a \"Kubernetes\" filter, "
         "while one that lists the word in a skills footer passes."),
        ("Generic \"rank these CVs\" AI is unexplainable",
         "Asking a language model to order candidates produces a result that "
         "cannot be traced to any specific line of any CV, and that may differ "
         "between two runs on identical input."),
        ("Employment decisions demand auditability",
         "Where automation touches who gets hired, an organisation needs to show "
         "what was considered, what it concluded, and who acted on it — a "
         "requirement no black-box ranking can satisfy."),
    ], cols=2, height=Inches(1.9))
    footer(slide, n)
    notes(slide, "Frame the gap in the existing landscape without overclaiming: "
                 "no invented statistics here — every point is a structural "
                 "property of the approach being criticised.\n\n"
                 "The fourth card is the one that motivates the whole design: "
                 "auditability is a hard requirement, not a feature.")
    return slide


def slide_problem(prs, n):
    slide, top = content_slide(
        prs, "1 · Problem statement",
        "The problem this project addresses",
        "Recruiters do not need candidates ranked. They need to know who meets "
        "the stated requirements, and on what evidence.")

    sh = box(slide, MARGIN, top, CONTENT_W, Inches(0.95), fill=ACCENT_SOFT,
             line=None)
    box(slide, MARGIN, top, Inches(0.06), Inches(0.95), fill=ACCENT,
        shape=MSO_SHAPE.RECTANGLE)
    tf = sh.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.32)
    tf.margin_right = Inches(0.3)
    line_of(tf, "How can a recruitment system decide whether a candidate is "
                "qualified, in a way that is evidence-backed, reproducible, "
                "and always attributable to a human decision-maker?",
            size=16, color=INK, bold=True, first=True, space_after=0)

    bullets(slide, top + Inches(1.3), [
        ("No traceability.  ", "A score of 74% tells a recruiter nothing about "
         "which requirement it failed on, or which line of the CV supports it. "
         "The number cannot be defended to a candidate or an auditor."),
        ("No reproducibility.  ", "If the score comes from a model sample, the "
         "same CV scored twice can yield two different outcomes — so the result "
         "is not a measurement of the candidate."),
        ("No accountability.  ", "When automation advances or rejects a "
         "candidate silently, there is no record of who decided, when, or why "
         "— and no way to override it on the record."),
        ("No separation of extraction from judgement.  ", "Systems that let a "
         "model both read the CV and decide the outcome make the two failure "
         "modes inseparable: a misreading and a misjudgement look identical."),
    ], size=14, gap=13)
    footer(slide, n)
    notes(slide, "Read the research question aloud — it is the sentence the "
                 "rest of the deck answers.\n\n"
                 "The four gaps map one-to-one onto the four design principles "
                 "on the methodology slide, and each is closed by a specific "
                 "implementation slide later. Say that explicitly: it is the "
                 "spine of the report.")
    return slide


def slide_objectives(prs, n):
    slide, top = content_slide(
        prs, "2 · Aim & objectives", "Aim and objectives")

    sh = box(slide, MARGIN, top, CONTENT_W, Inches(0.82), fill=INK, line=None)
    tf = sh.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.32)
    tf.margin_right = Inches(0.3)
    line_of(tf, "Aim   ", size=13, color=ACCENT_SOFT, bold=True, first=True,
            space_after=0)
    run(tf.paragraphs[0], "To build a working multi-tenant web application "
        "that qualifies candidates against explicit, weighted job requirements "
        "and presents every verdict with traceable evidence — while leaving "
        "the hiring decision with a recruiter.", size=13, color=WHITE)

    table(slide, top + Inches(1.12),
          ["", "Objective", "Delivered by"],
          [
              ["O1", "Parse a free-text job description into structured, weighted requirements",
               "AI Job Parser"],
              ["O2", "Extract structured, evidence-tagged candidate data from PDF/DOCX CVs",
               "AI CV Extractor"],
              ["O3", "Score candidates deterministically and reproducibly against requirements",
               "Rule-based engine"],
              ["O4", "Present every requirement verdict with its supporting CV excerpt",
               "Evidence-card UI"],
              ["O5", "Keep the decision with the recruiter and log every action with a reason",
               "Decisions + audit log"],
              ["O6", "Support multiple organisations with role-based access and self-service onboarding",
               "Org model + /signup"],
          ],
          widths=[0.055, 0.65, 0.295], row_h=Inches(0.55), size=12.5)
    footer(slide, n)
    notes(slide, "Six objectives, each traceable to a specific component — the "
                 "right-hand column is deliberately there so the Results slide "
                 "can be checked against it.\n\n"
                 "O6 is the newest: recruiter self-service onboarding replaced "
                 "an admin-only invite flow late in development.")
    return slide


def slide_scope(prs, n):
    slide, top = content_slide(
        prs, "3 · Scope", "Scope and delimitations",
        "Boundaries were set to keep the qualification problem itself in focus "
        "rather than rebuilding an entire HR suite.")
    two_col(slide, top,
            "In scope — implemented",
            [
                "Multi-tenant organisation model with role-based access",
                "Job and weighted-requirement management (AI-assisted or manual)",
                "Bulk CV upload with PDF/DOCX text extraction",
                "Deterministic, reproducible qualification scoring",
                "Per-requirement evidence presented in the UI",
                "Recruiter shortlist / review / reject decisions, all audited",
                "Public job board and applicant self-apply",
                "Email-OTP and optional password authentication",
                "Recruiter self-service signup with organisation creation",
                "Super-admin organisation and AI-provider administration",
            ],
            "Out of scope — deliberately excluded",
            [
                "Interview scheduling, offers, and onboarding",
                "Payroll or HRIS integration",
                "Background and reference checking",
                "Automated rejection without a recruiter action",
                "Blind screening (designed, not implemented — see Limitations)",
                "Queue-backed asynchronous processing at scale",
                "Per-organisation teammate invitations",
                "Automated regression test suite",
            ],
            height=Inches(4.35), item_size=11.5, item_space=6)
    footer(slide, n)
    notes(slide, "Be direct about the right-hand column — an examiner will "
                 "look for whether exclusions were chosen or simply missed.\n\n"
                 "Two entries are honest admissions rather than scope "
                 "decisions: blind screening is claimed in the landing-page "
                 "marketing copy but is not implemented, and there is no "
                 "automated test suite. Both reappear under Limitations.")
    return slide


def slide_methodology(prs, n):
    slide, top = content_slide(
        prs, "4 · Methodology", "Four design principles",
        "Each principle answers one of the four gaps in the problem statement, "
        "and each is enforced structurally rather than by convention.")
    cards(slide, top, [
        ("1 · Separate extraction from judgement",
         "The language model reads; it never decides. It converts unstructured "
         "CV text into a structured profile. A separate pure function turns "
         "that profile into a score. A misreading and a misjudgement are "
         "therefore diagnosable independently."),
        ("2 · Evidence-gating",
         "No requirement may be marked Confirmed unless the CV text supports "
         "it. Related-but-insufficient text must be marked Unclear with a "
         "reason; absent text is Not Found. Guessing is not an available "
         "output."),
        ("3 · Determinism & reproducibility",
         "The score is a weighted function of stored, structured data — not a "
         "model sample. It can be recomputed at any time from the persisted "
         "profile and requirements, and will return the same value."),
        ("4 · Human-in-the-loop by construction",
         "The system produces a qualification status, never an outcome. "
         "Shortlisting, review and rejection are recruiter actions, each "
         "written to an audit log with an actor and a reason."),
    ], cols=2, height=Inches(2.0))
    footer(slide, n)
    notes(slide, "This is the methodological core of the report — spend time "
                 "here.\n\n"
                 "Principle 1 is the design decision an examiner is most "
                 "likely to probe: why not just ask a model for the score? "
                 "Answer: because a sampled score is neither reproducible nor "
                 "auditable, and it collapses two distinct failure modes into "
                 "one unexplainable number.")
    return slide


def slide_architecture(prs, n):
    slide, top = content_slide(
        prs, "5 · Architecture", "System architecture",
        "A layered Next.js application; the AI layer is one replaceable "
        "service among several, not the centre of the system.")

    layers = [
        ("Presentation", "Next.js 14 App Router · React 18 server + client components · "
                         "Tailwind design tokens · Recharts visualisations", ACCENT_SOFT, ACCENT),
        ("Application / API", "26 route handlers · Zod request validation · NextAuth JWT sessions · "
                              "middleware role gate", SURFACE, ACCENT),
        ("Domain logic", "Rule-based scoring engine (pure) · qualification pipeline · "
                         "PDF/DOCX text extraction", SUCCESS_SOFT, SUCCESS),
        ("AI services", "Provider-agnostic OpenAI-compatible client · Job Parser · CV Extractor · "
                         "Explanation Service · JSON-mode with repair and retry", WARNING_SOFT, WARNING),
        ("Data", "PostgreSQL (Neon) via Prisma ORM · 15 models · 6 enums · 4 migrations",
         SURFACE, ACCENT),
        ("External services", "DropAphi — file storage, transactional email, hosted OTP delivery "
                              "and verification", CANVAS, MUTED),
    ]
    h = Inches(0.63)
    gap = Inches(0.12)
    for i, (name, detail, fill, accent) in enumerate(layers):
        t = top + i * (h + gap)
        sh = box(slide, MARGIN, t, CONTENT_W, h, fill=fill, line=BORDER)
        box(slide, MARGIN, t, Inches(0.055), h, fill=accent,
            shape=MSO_SHAPE.RECTANGLE)
        tf = sh.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.28)
        tf.margin_right = Inches(0.2)
        p = para(tf, True)
        p.space_after = Pt(0)
        run(p, name + "     ", size=13.5, color=accent, bold=True)
        run(p, detail, size=11.5, color=INK)
    footer(slide, n)
    notes(slide, "Read the layers bottom-up if asked about data flow, "
                 "top-down if asked about a user journey.\n\n"
                 "The AI-services layer is provider-agnostic on purpose: four "
                 "providers (AgentRouter, OpenAI, DeepSeek, OpenRouter) are "
                 "supported behind one OpenAI-compatible interface, selected "
                 "from a database row that falls back to environment "
                 "variables. Swapping model vendors requires no code change.")
    return slide


def slide_data_model(prs, n):
    slide, top = content_slide(
        prs, "5 · Data model", "Domain model",
        "Fifteen Prisma models. Every recruiter-facing query is scoped by "
        "organisation, which is what makes the system multi-tenant.")

    def node(left, t, w, h, label, sub, fill=SURFACE, accent=ACCENT):
        sh = box(slide, left, t, w, h, fill=fill, line=accent, line_pt=1.25)
        tf = sh.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.08)
        line_of(tf, label, size=12, color=INK, bold=True, first=True,
                space_after=2, align=PP_ALIGN.CENTER)
        if sub:
            line_of(tf, sub, size=9, color=MUTED, space_after=0,
                    align=PP_ALIGN.CENTER, font=MONO)
        return sh

    def arrow(left, t, label):
        a = box(slide, left, t + Inches(0.28), Inches(0.42), Inches(0.2),
                fill=ACCENT, shape=MSO_SHAPE.RIGHT_ARROW)
        a.line.fill.background()
        tf = text(slide, left - Inches(0.05), t + Inches(0.03), Inches(0.52),
                  Inches(0.22), align=PP_ALIGN.CENTER)
        line_of(tf, label, size=9, color=MUTED, first=True, space_after=0,
                align=PP_ALIGN.CENTER, font=MONO)

    w, h = Inches(2.35), Inches(0.78)
    step = w + Inches(0.42)

    r1 = top
    node(MARGIN, r1, w, h, "Organization", "tenant root")
    arrow(MARGIN + w, r1, "1–*")
    node(MARGIN + step, r1, w, h, "Job", "status · threshold")
    arrow(MARGIN + step + w, r1, "1–*")
    node(MARGIN + 2 * step, r1, w, h, "JobRequirement",
         "type · weight · mandatory")
    node(MARGIN + 3 * step, r1, w, h, "User", "role · organizationId",
         fill=ACCENT_SOFT)

    r2 = top + Inches(1.22)
    node(MARGIN, r2, w, h, "Candidate", "profile JSON")
    arrow(MARGIN + w, r2, "1–*")
    node(MARGIN + step, r2, w, h, "Application", "status · score")
    arrow(MARGIN + step + w, r2, "1–1")
    node(MARGIN + 2 * step, r2, w, h, "Evaluation", "score · status · reason",
         fill=SUCCESS_SOFT, accent=SUCCESS)
    arrow(MARGIN + 2 * step + w, r2, "1–*")
    node(MARGIN + 3 * step, r2, w, h, "EvaluationCriterion",
         "verdict · evidence", fill=SUCCESS_SOFT, accent=SUCCESS)

    r3 = top + Inches(2.44)
    node(MARGIN, r3, w, h, "CandidateSkill", "confidence · evidence")
    node(MARGIN + step, r3, w, h, "RecruiterDecision", "actor · reason",
         fill=WARNING_SOFT, accent=WARNING)
    node(MARGIN + 2 * step, r3, w, h, "AuditLog", "action · metadata",
         fill=WARNING_SOFT, accent=WARNING)
    node(MARGIN + 3 * step, r3, w, h, "AiSettings", "singleton row",
         fill=CANVAS, accent=MUTED)

    tf = text(slide, MARGIN, top + Inches(3.5), CONTENT_W, Inches(0.9))
    line_of(tf, "Users and Jobs both belong to an Organization, and every "
                "Application belongs to a Job — those three edges are what make "
                "the system multi-tenant. Each EvaluationCriterion points back "
                "at the JobRequirement it judged, so a score can always be "
                "decomposed into the requirements that produced it. "
                "RecruiterDecision and AuditLog hang off Application rather "
                "than Evaluation: human actions are recorded against the "
                "application, independently of whatever the AI concluded.",
            size=11.5, color=MUTED, first=True, space_after=0)
    line_of(tf, "Not shown: CVDocument / CVSection / CVVersion (candidate CV "
                "builder).", size=10.5, color=MUTED, italic=True,
            space_before=6, space_after=0)
    footer(slide, n)
    notes(slide, "Two structural points worth making if asked.\n\n"
                 "First, EvaluationCriterion is a join between an evaluation "
                 "and a requirement carrying a verdict and its evidence — that "
                 "single table is what makes the score decomposable.\n\n"
                 "Second, audit records attach to the Application, not the "
                 "Evaluation, so human actions survive a re-evaluation.")
    return slide


def slide_pipeline(prs, n):
    slide, top = content_slide(
        prs, "6 · Implementation", "The qualification pipeline",
        "One application at a time, written to be queue-ready: a single "
        "identifier in, no shared state.")
    flow(slide, top, [
        ("Upload", "PDF / DOCX\nto DropAphi"),
        ("Text extraction", "pdf-parse\nmammoth"),
        ("CV Extractor", "AI → structured\nevidence-tagged JSON"),
        ("Scoring engine", "deterministic\nrule-based"),
        ("Explanation", "AI → strengths,\ngaps, recommendation"),
        ("Persist", "Evaluation +\naudit log"),
    ], height=Inches(1.15))

    tf = text(slide, MARGIN, top + Inches(1.5), CONTENT_W, Inches(0.35))
    line_of(tf, "Application status transitions", size=13, color=INK,
            bold=True, first=True, space_after=0)

    states = [("UPLOADED", MUTED), ("PROCESSING", MUTED), ("EXTRACTED", ACCENT),
              ("EVALUATING", ACCENT), ("QUALIFIED", SUCCESS),
              ("REVIEW_REQUIRED", WARNING), ("NOT_QUALIFIED", DANGER)]
    x = MARGIN
    t = top + Inches(1.92)
    for i, (state, color) in enumerate(states):
        w = Inches(1.62)
        sh = box(slide, x, t, w, Inches(0.42), fill=SURFACE, line=color)
        tf2 = sh.text_frame
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        line_of(tf2, state, size=9.5, color=color, bold=True, first=True,
                space_after=0, align=PP_ALIGN.CENTER, font=MONO)
        x += w
        if i < len(states) - 1:
            sep = "→" if i < 4 else "/"
            tf2 = text(slide, x, t + Inches(0.08), Inches(0.28), Inches(0.3),
                       align=PP_ALIGN.CENTER)
            line_of(tf2, sep, size=12, color=MUTED, first=True, space_after=0,
                    align=PP_ALIGN.CENTER)
            x += Inches(0.28)

    bullets(slide, top + Inches(2.62), [
        ("Failure isolation.  ", "A failed extraction returns the application "
         "to UPLOADED and surfaces the error — a half-processed record is never "
         "shown as a result. A failed explanation does not invalidate the "
         "score: the verdict does not depend on the prose describing it."),
        ("Re-evaluation is idempotent.  ", "Evaluations are upserted and their "
         "criteria replaced, so requirements can be corrected and every "
         "application rescored without duplicating records."),
        ("Every completed run is audited.  ", "A system-actor entry records the "
         "resulting score and status, so an AI verdict is as traceable as a "
         "human one."),
    ], size=12, gap=8)
    footer(slide, n)
    notes(slide, "The honest architectural caveat belongs here: this pipeline "
                 "currently runs synchronously inside the HTTP request. It is "
                 "written to be queue-ready — one application id in, no shared "
                 "state — but the queue itself is future work, and at the "
                 "hundreds-of-CVs scale in the brief it would be required.\n\n"
                 "Note the deliberate asymmetry in failure handling: an "
                 "extraction failure invalidates the result, an explanation "
                 "failure does not.")
    return slide


def slide_scoring(prs, n):
    slide, top = content_slide(
        prs, "6 · Implementation", "The scoring engine",
        "A pure function: requirements plus an extracted profile in, a score "
        "and status out. No model call, no randomness, no hidden state.")

    sh = box(slide, MARGIN, top, CONTENT_W, Inches(0.62), fill=INK, line=None)
    tf = sh.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.32)
    line_of(tf, "score  =  round( Σ  criterion_score × ( weight / Σ weights ) )",
            size=15, color=WHITE, bold=True, first=True, space_after=0,
            font=MONO)

    left_w = int(CONTENT_W * 0.47)
    tf = text(slide, MARGIN, top + Inches(0.95), left_w, Inches(0.3))
    line_of(tf, "Criterion verdicts", size=13, color=INK, bold=True,
            first=True, space_after=0)
    verdicts = [("Confirmed", "100", SUCCESS, SUCCESS_SOFT),
                ("Unclear", "50", WARNING, WARNING_SOFT),
                ("Not Found", "0", DANGER, DANGER_SOFT)]
    for i, (label, val, color, fill) in enumerate(verdicts):
        sh = box(slide, MARGIN, top + Inches(1.32) + i * Inches(0.52),
                 left_w, Inches(0.44), fill=fill, line=None)
        tf = sh.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.22)
        p = para(tf, True)
        p.space_after = Pt(0)
        run(p, label, size=12.5, color=color, bold=True)
        run(p, "   scores   ", size=11.5, color=MUTED)
        run(p, val, size=12.5, color=color, bold=True, font=MONO)

    tf = text(slide, MARGIN, top + Inches(3.02), left_w, Inches(2.0))
    line_of(tf, "Matching rules", size=13, color=INK, bold=True, first=True,
            space_after=8)
    for s in [
        "Experience requirements (\"5+ years\") are matched numerically: at or "
        "above the requirement → Confirmed; at or above 60% of it → Unclear; "
        "below that → Not Found.",
        "Skill, tool, certification and education requirements are matched by "
        "normalised name against the extracted skills list, inheriting the "
        "evidence-gated verdict the extractor assigned.",
    ]:
        p = para(tf)
        p.space_after = Pt(7)
        run(p, "■   ", size=10, color=ACCENT, bold=True)
        run(p, s, size=11.5, color=INK)

    right_x = MARGIN + left_w + Inches(0.4)
    right_w = CONTENT_W - left_w - Inches(0.4)
    tf = text(slide, right_x, top + Inches(0.95), right_w, Inches(0.3))
    line_of(tf, "Status bands", size=13, color=INK, bold=True, first=True,
            space_after=0)

    bands = [("Strong Match", "≥ 85 and no unclear mandatory", SUCCESS),
             ("Qualified", "≥ 70", SUCCESS),
             ("Needs Review", "55 – 69, or gated below", WARNING),
             ("Not Qualified", "< 55", DANGER)]
    for i, (label, cond, color) in enumerate(bands):
        t = top + Inches(1.32) + i * Inches(0.52)
        sh = box(slide, right_x, t, right_w, Inches(0.44), fill=SURFACE,
                 line=BORDER)
        box(slide, right_x, t, Inches(0.05), Inches(0.44), fill=color,
            shape=MSO_SHAPE.RECTANGLE)
        tf = sh.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.22)
        p = para(tf, True)
        p.space_after = Pt(0)
        run(p, label, size=12.5, color=color, bold=True)
        run(p, "     " + cond, size=11, color=MUTED, font=MONO)

    sh = box(slide, right_x, top + Inches(3.02), right_w, Inches(1.45),
             fill=ACCENT_SOFT, line=None)
    box(slide, right_x, top + Inches(3.02), Inches(0.06), Inches(1.45),
        fill=ACCENT, shape=MSO_SHAPE.RECTANGLE)
    tf = sh.text_frame
    tf.margin_left = Inches(0.28)
    tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.16)
    line_of(tf, "Two overrides on top of the weighted score", size=12.5,
            color=ACCENT, bold=True, first=True, space_after=6)
    line_of(tf, "Mandatory Requirement Override — any mandatory requirement "
                "marked Not Found caps the outcome at Needs Review, however "
                "high the weighted score.", size=11, color=INK, space_after=5)
    line_of(tf, "Recruiter threshold — a per-job threshold can demote "
                "Qualified to Needs Review, letting a team tighten the bar "
                "without touching the engine.", size=11, color=INK,
            space_after=0)
    footer(slide, n)
    notes(slide, "The most technical slide; it directly answers objective O3 "
                 "and the reproducibility gap.\n\n"
                 "The mandatory override is the part to emphasise — it is why "
                 "a high average cannot paper over a missing hard requirement. "
                 "A candidate scoring 92% who lacks a mandatory qualification "
                 "surfaces as Needs Review with a stated reason, not as a "
                 "Strong Match.\n\n"
                 "If asked why 55/70/85 and 60%: they are tuned defaults, and "
                 "because the engine is a pure function they can be changed "
                 "and every historical application rescored deterministically.")
    return slide


def slide_extraction(prs, n):
    slide, top = content_slide(
        prs, "6 · Implementation", "Evidence-gated extraction",
        "The extractor is constrained by an explicit contract; the engineering "
        "around it exists because real gateways and models misbehave.")

    left_w = int(CONTENT_W * 0.5) - Inches(0.2)
    sh = box(slide, MARGIN, top, left_w, Inches(3.62), fill=SURFACE,
             line=BORDER)
    tf = sh.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.22)
    line_of(tf, "The extraction contract", size=13.5, color=ACCENT, bold=True,
            first=True, space_after=10)
    for s in [
        "\"Confirmed\" only when the CV text directly names the skill, or "
        "clearly implies it through a specific named tool or technology.",
        "A broader category without the specific requirement — \"cloud "
        "experience\" against \"AWS\" — must be marked \"unclear\", with the "
        "reason recorded as evidence.",
        "Evidence must be a short direct excerpt or close paraphrase from the "
        "CV. Invented quotations are prohibited.",
        "No skill, employer, role or qualification may be returned that is not "
        "present in the text.",
        "Age, gender, photograph, marital status, nationality, race and "
        "religion are excluded at extraction — so they are structurally absent "
        "from everything the scoring engine can see.",
    ]:
        p = para(tf)
        p.space_after = Pt(7)
        run(p, "■   ", size=9, color=ACCENT, bold=True)
        run(p, s, size=11, color=INK)

    right_x = MARGIN + left_w + Inches(0.4)
    right_w = CONTENT_W - left_w - Inches(0.4)
    sh = box(slide, right_x, top, right_w, Inches(3.62), fill=CANVAS,
             line=BORDER)
    tf = sh.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.22)
    line_of(tf, "Making a real provider behave", size=13.5, color=WARNING,
            bold=True, first=True, space_after=10)
    for s in [
        "JSON mode is requested, but gateway-hosted models do not reliably "
        "comply — responses are repaired before parsing rather than trusted.",
        "Reasoning models can spend the whole token budget on hidden reasoning "
        "and return empty content; this is detected and reported as a distinct "
        "failure rather than an empty profile.",
        "A misconfigured gateway base URL can return an HTML console page with "
        "HTTP 200; URL normalisation removes the class of failure entirely.",
        "One gateway rejects unrecognised client User-Agents with a 401 that is "
        "indistinguishable from a bad key — a known-good header is sent.",
        "Transient failures are retried with backoff; CV text is truncated to a "
        "token budget before the call.",
    ]:
        p = para(tf)
        p.space_after = Pt(7)
        run(p, "■   ", size=9, color=WARNING, bold=True)
        run(p, s, size=11, color=INK)

    sh = box(slide, MARGIN, top + Inches(3.85), CONTENT_W, Inches(0.72),
             fill=ACCENT_SOFT, line=None)
    box(slide, MARGIN, top + Inches(3.85), Inches(0.06), Inches(0.72),
        fill=ACCENT, shape=MSO_SHAPE.RECTANGLE)
    tf = sh.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.25)
    line_of(tf, "Why this matters for fairness: the scoring engine receives "
                "only skills, experience, education and certifications. "
                "Protected attributes are not weighted at zero — they are "
                "never in the data structure the engine reads.",
            size=12, color=INK, bold=True, first=True, space_after=0)
    footer(slide, n)
    notes(slide, "The left column answers objectives O2 and O4; the right "
                 "column is the part that shows engineering maturity, because "
                 "every item is a failure that actually occurred and was "
                 "diagnosed during development.\n\n"
                 "The band at the bottom is the fairness argument, and it is "
                 "worth stating precisely: exclusion is structural, not a "
                 "weighting choice. Do not overclaim beyond it — bias can "
                 "still enter through proxies in the CV text, which is "
                 "acknowledged under Limitations.")
    return slide


def slide_auth(prs, n):
    slide, top = content_slide(
        prs, "6 · Implementation", "Authentication and onboarding",
        "Three sign-in methods and two routes into a recruiter account — the "
        "second added to remove the admin bottleneck.")

    tf = text(slide, MARGIN, top, CONTENT_W, Inches(0.3))
    line_of(tf, "Path A — administrator-provisioned (original flow)", size=13,
            color=MUTED, bold=True, first=True, space_after=0)
    flow(slide, top + Inches(0.36), [
        ("Super admin creates org", "at /admin"),
        ("Invite email sent", "via DropAphi"),
        ("Recruiter signs in", "emailed code"),
        ("Lands in org dashboard", "role pre-assigned"),
    ], height=Inches(0.78))

    tf = text(slide, MARGIN, top + Inches(1.42), CONTENT_W, Inches(0.3))
    line_of(tf, "Path B — recruiter self-service (added)", size=13,
            color=ACCENT, bold=True, first=True, space_after=0)
    flow(slide, top + Inches(1.78), [
        ("Signup from hero", "/signup"),
        ("Code emailed", "nothing persisted yet"),
        ("Code verified", "org + recruiter created\nin one transaction"),
        ("Straight to job creation", "session already started"),
    ], height=Inches(0.78))

    cards(slide, top + Inches(2.92), [
        ("Nothing is written before the code verifies",
         "The organisation and account are created only after the emailed code "
         "is confirmed, so an unverified address can never come to own an "
         "organisation."),
        ("Roles are never taken from the client",
         "The role is resolved server-side from the stored record after "
         "authentication, whichever method was used. Routes are gated in "
         "middleware and re-checked per endpoint."),
        ("Passwords optional, never displaced",
         "Only a bcrypt hash is stored and the email-code route always remains "
         "available. Holding a code does not replace an existing password."),
        ("Disclosure is bounded",
         "A signup form must refuse duplicates, so it reveals only that an "
         "address already belongs to an organisation — never the role, the org, "
         "or whether a password is set."),
    ], cols=4, height=Inches(1.7), heading_size=12, body_size=10.5)
    footer(slide, n)
    notes(slide, "Path B is the most recent piece of work and a good "
                 "demonstration target: the original design required a super "
                 "admin to create every organisation and email an invite, "
                 "which made the platform unusable without operator "
                 "involvement.\n\n"
                 "The security reasoning in the first card is the point to "
                 "make: deferring all writes until after verification is what "
                 "prevents unverified-email account creation, and it is why "
                 "the form data travels with the verification code rather than "
                 "being stored as a pending record first.")
    return slide


def slide_access(prs, n):
    slide, top = content_slide(
        prs, "6 · Implementation", "Roles and multi-tenancy",
        "Four application roles plus a platform tier. Organisation scoping is "
        "applied in the query, not in the interface.")
    table(slide, top,
          ["Role", "Can do", "Scope"],
          [
              ["SUPERADMIN", "Create organisations, invite their first recruiter, "
               "configure the platform AI provider", "Platform-wide"],
              ["RECRUITER", "Create jobs and requirements, upload CVs, shortlist / "
               "request review / reject, generate shortlist emails", "Own organisation"],
              ["HIRING_MANAGER", "Review qualified candidates and their evidence "
               "within the organisation", "Own organisation"],
              ["APPLICANT", "Browse the public board, self-apply, view own scores "
               "and evidence", "Own applications"],
          ],
          widths=[0.19, 0.585, 0.225], row_h=Inches(0.62), size=12)

    bullets(slide, top + Inches(3.25), [
        ("Two-stage enforcement.  ", "Middleware gates whole route groups by "
         "role; each endpoint independently re-derives the caller's user and "
         "organisation from the session before touching data."),
        ("Tenancy in the query.  ", "Organisation-owned records are always "
         "fetched by organizationId rather than filtered after loading, so a "
         "missing check fails closed and returns nothing."),
    ], size=12.5, gap=9)
    footer(slide, n)
    notes(slide, "If asked about horizontal privilege escalation: tenancy is "
                 "enforced in the query itself, so the failure mode of a "
                 "forgotten check is an empty result rather than another "
                 "organisation's data.\n\n"
                 "Applicant self-apply shares the same authentication as "
                 "recruiters — only the stored role differs, which is what "
                 "keeps one auth flow serving every user type.")
    return slide


def slide_results(prs, n):
    slide, top = content_slide(
        prs, "7 · Results", "What was delivered",
        "A working application, not a prototype of one — every objective has a "
        "running implementation behind it.")
    tiles(slide, top, [
        ("11", "application pages"),
        ("26", "API endpoints"),
        ("22", "React components"),
        ("15", "data models"),
        ("6", "domain enums"),
        ("4", "AI providers supported"),
        ("~6.7k", "lines of TypeScript"),
        ("0", "AI-made hiring decisions"),
    ], cols=4, height=Inches(1.1))

    cards(slide, top + Inches(2.68), [
        ("O1 · Job → requirements", "AI-generated or manual, editable either way"),
        ("O2 · Evidence-tagged extraction", "PDF/DOCX, gated at the prompt level"),
        ("O3 · Deterministic scoring", "Pure function, recomputable at any time"),
        ("O4 · Evidence per verdict", "Evidence cards, one per requirement"),
        ("O5 · Human decision, logged", "Decisions plus a full audit trail"),
        ("O6 · Multi-tenant + self-signup", "Org model, RBAC, recruiter signup"),
    ], cols=3, height=Inches(0.8), gap=Inches(0.2), heading_size=11.5,
        body_size=10, accent=SUCCESS, fill=SUCCESS_SOFT)
    footer(slide, n)
    notes(slide, "Counts are measured from the repository, not estimated — the "
                 "script that builds this deck reports them so the numbers "
                 "cannot drift from the code.\n\n"
                 "The last tile is a deliberate design claim rather than a "
                 "metric: the system produces qualification statuses, and no "
                 "code path advances or rejects a candidate without a "
                 "recruiter action.")
    return slide


def slide_output(prs, n):
    slide, top = content_slide(
        prs, "7 · Results", "What the recruiter actually sees",
        "The interface is built around one signature element — the evidence "
        "card — repeated everywhere a conclusion is drawn.")
    cards(slide, top, [
        ("Organisation dashboard",
         "Qualified, Needs Review and Not Qualified counts per job across the "
         "whole organisation, so the funnel is visible before any individual "
         "CV is opened."),
        ("Candidate list per job",
         "Every application with its score and qualification status, ordered "
         "for triage rather than ranked for a decision."),
        ("Per-requirement evidence",
         "Each requirement marked Confirmed, Unclear or Not Found, with the "
         "supporting excerpt from the CV shown alongside its confidence."),
        ("Generated explanation",
         "Strengths, gaps and a recommendation in prose — derived from the "
         "already-computed score, never the source of it."),
        ("Recruiter actions, all recorded",
         "Shortlist, request review, reject and free-text notes, each written "
         "to the audit trail with an actor and a reason."),
        ("Applicant transparency",
         "Applicants see their own scores across the roles they applied to — "
         "the same evidence the recruiter sees, not a bare rejection."),
    ], cols=3, height=Inches(1.85), body_size=11.5)
    footer(slide, n)
    notes(slide, "If there is a live demo, this is the slide to demo from — "
                 "dashboard, then one job, then one candidate, then the "
                 "evidence.\n\n"
                 "The last card matters for the ethics discussion: applicants "
                 "get the same evidence rather than an unexplained outcome, "
                 "which is unusual in screening tools.")
    return slide


def slide_evaluation(prs, n):
    slide, top = content_slide(
        prs, "8 · Evaluation", "Verification and its limits",
        "What was actually established about correctness — and, as clearly, "
        "what was not.")
    two_col(slide, top,
            "Verified",
            [
                "End-to-end manual runs of the full pipeline: job creation, "
                "requirement generation, CV upload, scoring, explanation, "
                "recruiter decision",
                "TypeScript strict typecheck passes across the codebase",
                "Production build compiles — all 11 pages and 26 endpoints",
                "AI provider connectivity verified by a dedicated check script",
                "Signup validation exercised directly against the running API: "
                "weak password, short organisation name and duplicate address "
                "each rejected with the intended status code",
                "Known provider failure modes reproduced and handled: HTML-200 "
                "gateway responses, empty reasoning-model content, truncated "
                "JSON, User-Agent 401s",
            ],
            "Not established",
            [
                "No automated regression test suite exists — the scoring "
                "engine is a pure function and unit-testable by construction, "
                "but those tests are not yet written",
                "No measured accuracy benchmark against human screeners on a "
                "labelled CV corpus",
                "No load testing; the synchronous pipeline is untested at the "
                "hundreds-of-CVs scale in the brief",
                "No formal inter-rater study of whether the Confirmed / "
                "Unclear / Not Found verdicts match human judgement",
                "Extraction quality not quantified across CV formats, "
                "languages, or scanned documents",
            ],
            left_accent=SUCCESS, right_accent=WARNING,
            left_fill=SUCCESS_SOFT, right_fill=WARNING_SOFT,
            height=Inches(4.35), item_size=11, item_space=6)
    footer(slide, n)
    notes(slide, "Do not soften the right-hand column. Naming the gap between "
                 "'it works when I run it' and 'its accuracy is measured' is "
                 "what an examiner is looking for, and claiming otherwise is "
                 "the fastest way to lose credibility.\n\n"
                 "The strongest thing to say: because scoring is a pure "
                 "function of stored data, a benchmark can be run "
                 "retrospectively over any labelled corpus without touching "
                 "the engine. The design does not block the evaluation — it "
                 "simply has not been done yet.")
    return slide


def slide_limitations(prs, n):
    slide, top = content_slide(
        prs, "9 · Limitations", "Known limitations",
        "Stated as they stand in the current implementation.")
    cards(slide, top, [
        ("Synchronous processing",
         "CV processing runs inside the HTTP request. The pipeline is written "
         "queue-ready, but at hundreds of CVs per vacancy an asynchronous "
         "worker is required, not optional."),
        ("Literal requirement matching",
         "Requirements are matched by normalised name and substring, so "
         "semantically equivalent phrasings — \"K8s\" against \"Kubernetes\" — "
         "can be missed and reported as Not Found."),
        ("Blind screening is claimed but not built",
         "The landing page advertises blind screening. No such capability "
         "exists in the codebase; the marketing copy overstates the system and "
         "should be corrected or the feature implemented."),
        ("Extraction bounded by input quality",
         "A scanned PDF with no text layer yields nothing to extract, as there "
         "is no OCR stage. Fairness also depends on the model honouring the "
         "exclusion instruction, which is not independently enforced."),
        ("Rate limiting is per-instance",
         "The one-code-per-30-seconds throttle is held in process memory, so "
         "it does not hold across a multi-instance deployment."),
        ("Single administrative tier",
         "There is one platform-wide super admin and no per-organisation admin "
         "role, so a recruiter cannot yet invite a teammate into their own "
         "organisation."),
    ], cols=3, height=Inches(1.9), body_size=11.5)
    footer(slide, n)
    notes(slide, "Volunteering the blind-screening discrepancy is deliberate. "
                 "It was found by checking the marketing claims against the "
                 "code, and reporting it is more defensible than letting an "
                 "examiner discover an unimplemented advertised feature.\n\n"
                 "The proxy-bias point under extraction is the honest limit of "
                 "the fairness claim: excluding protected attributes from the "
                 "data structure does not eliminate proxies present in the CV "
                 "text itself.")
    return slide


def slide_future(prs, n):
    slide, top = content_slide(
        prs, "9 · Future work", "Future work",
        "Each item follows directly from a limitation on the previous slide.")
    table(slide, top,
          ["Priority", "Work", "Addresses"],
          [
              ["1", "Queue-backed asynchronous processing with per-file progress in the UI",
               "Synchronous pipeline"],
              ["2", "Automated test suite — unit tests for the scoring engine, "
               "integration tests for the pipeline", "No regression safety net"],
              ["3", "Embedding-based semantic requirement matching, with the "
               "evidence excerpt retained", "Literal matching"],
              ["4", "Implement blind screening — identity masked until shortlist "
               "— or remove the claim", "Overstated capability"],
              ["5", "Per-organisation admin role and teammate invitations",
               "Single admin tier"],
              ["6", "Accuracy benchmark against human screeners on a labelled corpus",
               "Unmeasured accuracy"],
              ["7", "Redis-backed rate limiting and multi-instance deployment",
               "Per-instance throttle"],
              ["8", "Fairness auditing dashboard over historical decisions and overrides",
               "Unmonitored bias risk"],
          ],
          widths=[0.11, 0.62, 0.27], row_h=Inches(0.53), size=12)
    footer(slide, n)
    notes(slide, "Ordered by what unblocks the most: the queue makes the "
                 "system usable at the scale in the brief, and the test suite "
                 "makes every later change safe.\n\n"
                 "Items 4 and 6 are integrity work rather than features — "
                 "aligning the claims with the build, and measuring what the "
                 "system actually achieves.")
    return slide


def slide_conclusion(prs, n):
    slide = blank(prs)
    paint(slide, INK)
    box(slide, Inches(0), Inches(0), Inches(0.28), SLIDE_H, fill=ACCENT,
        shape=MSO_SHAPE.RECTANGLE)

    tf = text(slide, Inches(1.15), Inches(0.85), Inches(10.8), Inches(0.35))
    line_of(tf, "10 · CONCLUSION", size=12, color=ACCENT_SOFT, bold=True,
            first=True, space_after=0)

    tf = text(slide, Inches(1.15), Inches(1.32), Inches(10.8), Inches(0.6))
    line_of(tf, "Qualification, not ranking", size=34, color=WHITE, bold=True,
            first=True, space_after=0)

    box(slide, Inches(1.15), Inches(2.12), Inches(1.5), Inches(0.055),
        fill=ACCENT, shape=MSO_SHAPE.RECTANGLE)

    items = [
        ("The AI reads; it never decides.", "Extraction and judgement are "
         "separate components, so a misreading and a misjudgement can be "
         "diagnosed apart from one another."),
        ("Every score is decomposable and reproducible.", "A deterministic "
         "weighted function over stored structured data, recomputable at any "
         "time, with each requirement traceable to the CV text behind it."),
        ("Mandatory requirements cannot be averaged away.", "An unmet "
         "mandatory requirement caps the outcome regardless of the weighted "
         "score, with the reason stated."),
        ("A human decides, on the record.", "The system produces a status; "
         "shortlisting, review and rejection are recruiter actions, each "
         "audited with an actor and a reason."),
        ("Recruiters onboard themselves.", "Self-service signup creates the "
         "organisation only after the email is verified, removing the "
         "administrator bottleneck without weakening account security."),
    ]
    tf = text(slide, Inches(1.15), Inches(2.52), Inches(10.6), Inches(3.4))
    first = True
    for head, body in items:
        p = para(tf, first)
        first = False
        p.space_after = Pt(11)
        run(p, "■   ", size=11, color=ACCENT_SOFT, bold=True)
        run(p, head + "  ", size=14, color=WHITE, bold=True)
        run(p, body, size=13, color=RGBColor.from_string("A9B6C4"))

    tf = text(slide, Inches(1.15), Inches(6.35), Inches(10.6), Inches(0.6))
    line_of(tf, "Qualify first. Hire smarter.", size=19, color=ACCENT_SOFT,
            bold=True, italic=True, first=True, space_after=4)
    line_of(tf, "Thank you — questions welcome.", size=12,
            color=RGBColor.from_string("8A97A6"), space_after=0)

    notes(slide, "Close on the distinction the whole project rests on: the "
                 "system qualifies candidates against stated requirements "
                 "rather than ranking them against each other.\n\n"
                 "Likely questions: why not let the model score (see "
                 "reproducibility and auditability); how is fairness enforced "
                 "(structural exclusion at extraction, with the proxy caveat); "
                 "does it scale (not yet — the queue is future work item 1).")
    return slide


SLIDE_BUILDERS = [
    slide_contents, slide_background, slide_problem, slide_objectives,
    slide_scope, slide_methodology, slide_architecture, slide_data_model,
    slide_pipeline, slide_scoring, slide_extraction, slide_auth,
    slide_access, slide_results, slide_output, slide_evaluation,
    slide_limitations, slide_future,
]


def build(path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    for i, fn in enumerate(SLIDE_BUILDERS, start=2):
        fn(prs, i)
    total = len(SLIDE_BUILDERS) + 2
    slide_conclusion(prs, total)

    prs.save(path)
    return total


if __name__ == "__main__":
    if len(sys.argv) > 1:
        out = sys.argv[1]
    else:
        desktop = os.path.join(os.path.expanduser("~"), "Desktop")
        out = os.path.join(desktop, "TalentBridge-AI-Overview.pptx")

    count = build(out)
    print(f"Wrote {out} ({count} slides)")
